import os
import json
import uuid
import subprocess
import tempfile
from google import genai
from google.genai import types

GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")

def extract_text_from_file(file_storage) -> str:
    """Extract plain text from a .txt, .pdf, .docx, or .pptx upload."""
    filename = file_storage.filename.lower()

    if filename.endswith(".txt"):
        raw = file_storage.read()
        try:
            return raw.decode("utf-8")
        except UnicodeDecodeError:
            return raw.decode("latin-1")

    with tempfile.NamedTemporaryFile(suffix=os.path.splitext(filename)[1], delete=False) as tmp:
        file_storage.save(tmp.name)
        tmp_path = tmp.name

    try:
        if filename.endswith(".pdf"):
            try:
                import fitz  # PyMuPDF
            except ImportError:
                raise RuntimeError("PyMuPDF not installed. Run: pip install PyMuPDF")
            doc = fitz.open(tmp_path)
            text = "\n".join(page.get_text("text") for page in doc)
            doc.close()

        elif filename.endswith(".docx"):
            try:
                import docx
            except ImportError:
                raise RuntimeError("python-docx not installed. Run: pip install python-docx")
            doc = docx.Document(tmp_path)
            text = "\n".join([para.text for para in doc.paragraphs])

        elif filename.endswith(".pptx"):
            try:
                from pptx import Presentation
            except ImportError:
                raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
            prs = Presentation(tmp_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            text = "\n".join(text_runs)

        else:
            raise ValueError("Unsupported file type. Please upload .txt, .pdf, .docx, or .pptx.")

    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

    text = text.encode("utf-8", errors="ignore").decode("utf-8")
    return text


def generate_slide_data_with_gemini(text: str, mode: str, slide_topics: list, slide_count: int, language: str = "English", grade: str = "Primary School") -> list:
    """Call Gemini AI to generate structured slide data safely."""
    client = genai.Client(api_key=GEMINI_API_KEY)

    persona_setup = f"""You are an expert primary school educator familiar with the Maharashtra State Board syllabus.
Generate an engaging presentation tailored for {grade} students.
CRITICAL: The ENTIRE presentation (all titles and bullet points) MUST be written in {language}.
DO NOT use markdown formatting like **bold** or *italics*. Provide plain text only."""

    document_context = f"\nAnalyze the document below:\n\"\"\"\n{text[:6000]}\n\"\"\"" if text else "\nUse your extensive educational knowledge to accurately generate the content based on the provided topics."

    if mode == "manual" and slide_topics:
        topic_list = "\n".join(f"- {t}" for t in slide_topics)
        prompt = f"""{persona_setup}
{document_context}

Generate exactly {len(slide_topics)} slides for these specific topics:
{topic_list}

For each topic, create a slide with:
- A concise title (max 8 words)
- 3 to 5 bullet points (each max 15 words)
- layout: one of "title", "content", "two_column", "stat"

Respond ONLY with valid JSON array like:
[
  {{
    "title": "Slide Title",
    "bullets": ["Point one", "Point two"],
    "layout": "content"
  }}
]
No markdown, no preamble, only the JSON array."""
    else:
        prompt = f"""{persona_setup}
{document_context}

Generate a structured presentation with {slide_count} slides based on the document provided.

Include:
- 1 title slide (layout: "title") with a subtitle bullet
- {slide_count - 2} content slides (layout: "content" or "two_column" or "stat")
- 1 summary/conclusion slide (layout: "content")

Respond ONLY with valid JSON array like:
[
  {{
    "title": "Slide Title",
    "bullets": ["Point one"],
    "layout": "content"
  }}
]
No markdown, no preamble, only the JSON array."""

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt,
        config=types.GenerateContentConfig(
            thinking_config=types.ThinkingConfig(thinking_budget=0)
        )
    )

    raw = response.text.strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    raw = raw.strip()

    try:
        slides = json.loads(raw)
    except Exception as e:
        print(f"❌ Gemini returned invalid JSON: {raw}")
        raise ValueError("Failed to parse Gemini response as JSON.")

    if isinstance(slides, dict):
        for key, value in slides.items():
            if isinstance(value, list):
                slides = value
                break

    if not isinstance(slides, list):
        raise ValueError("Gemini did not return a valid list of slides.")

    # 🚀 FIXED: Clean markdown asterisks from the generated text!
    for slide in slides:
        if 'title' in slide and isinstance(slide['title'], str):
            slide['title'] = slide['title'].replace('**', '').replace('*', '')
        if 'bullets' in slide and isinstance(slide['bullets'], list):
            slide['bullets'] = [str(b).replace('**', '').replace('*', '') for b in slide['bullets']]

    return slides


def build_pptx_from_slides(slides: list, output_path: str) -> None:
    """Use pptxgenjs via Node.js to build the .pptx file safely from slide data."""

    PALETTE = {
        "primary": "065A82",
        "secondary": "1C7293",
        "accent": "21295C",
        "white": "FFFFFF",
        "light_bg": "EEF6FB",
        "text_dark": "1A1A2E",
        "text_muted": "4A6572",
    }

    js_slides = json.dumps(slides, ensure_ascii=True)
    palette_js = json.dumps(PALETTE)
    output_path_js = output_path.replace("\\", "/")

    node_script = f"""
const pptxgen = require("pptxgenjs")

const slides = {js_slides};
const C = {palette_js};

const pres = new pptxgen();
pres.layout = "LAYOUT_16x9";
pres.title = "Shiksha Sahayak Presentation";
pres.author = "Shiksha Sahayak AI";

function makeShadow() {{
  return {{ type: "outer", blur: 8, offset: 3, angle: 135, color: "000000", opacity: 0.12 }};
}}

slides.forEach((slide, idx) => {{
  const s = pres.addSlide();
  const layout = slide.layout || "content";

  const bullets = Array.isArray(slide.bullets) ? slide.bullets : [];
  const safeTitle = slide.title ? String(slide.title) : "Topic";

  if (layout === "title") {{
    s.background = {{ color: C.accent }};
    s.addShape(pres.shapes.RECTANGLE, {{ x: 7.5, y: 0, w: 2.5, h: 5.625, fill: {{ color: C.primary, transparency: 40 }}, line: {{ color: C.primary, transparency: 40 }} }});
    s.addShape(pres.shapes.RECTANGLE, {{ x: 8.5, y: 0, w: 1.5, h: 5.625, fill: {{ color: C.secondary, transparency: 30 }}, line: {{ color: C.secondary, transparency: 30 }} }});
    s.addShape(pres.shapes.ROUNDED_RECTANGLE, {{ x: 0.5, y: 0.35, w: 0.5, h: 0.3, fill: {{ color: C.secondary }}, line: {{ color: C.secondary }}, rectRadius: 0.05 }});
    s.addText(String(idx + 1), {{ x: 0.5, y: 0.35, w: 0.5, h: 0.3, fontSize: 10, color: C.white, align: "center", valign: "middle", bold: true, margin: 0 }});
    s.addText(safeTitle, {{ x: 0.6, y: 1.5, w: 6.8, h: 1.5, fontSize: 40, fontFace: "Cambria", bold: true, color: C.white, align: "left", valign: "middle" }});

    if (bullets.length > 0) {{
        s.addText(String(bullets[0]), {{ x: 0.6, y: 3.2, w: 6.8, h: 0.8, fontSize: 16, fontFace: "Calibri", color: "CADCFC", align: "left", italic: true }});
    }}

    s.addShape(pres.shapes.RECTANGLE, {{ x: 0, y: 5.2, w: 10, h: 0.425, fill: {{ color: C.primary }}, line: {{ color: C.primary }} }});
    s.addText("Shiksha Sahayak · AI-Powered Education", {{ x: 0.5, y: 5.2, w: 9, h: 0.425, fontSize: 10, color: C.white, align: "left", valign: "middle", margin: 0 }});

  }} else if (layout === "stat") {{
    s.background = {{ color: C.light_bg }};
    s.addShape(pres.shapes.RECTANGLE, {{ x: 0, y: 0, w: 10, h: 1.0, fill: {{ color: C.primary }}, line: {{ color: C.primary }} }});
    s.addText(safeTitle, {{ x: 0.5, y: 0, w: 9, h: 1.0, fontSize: 22, fontFace: "Cambria", bold: true, color: C.white, align: "left", valign: "middle", margin: 0 }});

    const cols = Math.min(bullets.length, 3);
    if (cols > 0) {{
        const cardW = (10 - 1.0) / cols;
        bullets.slice(0, 3).forEach((b, i) => {{
          const cx = 0.5 + i * (cardW + 0.1);
          s.addShape(pres.shapes.RECTANGLE, {{ x: cx, y: 1.3, w: cardW - 0.1, h: 2.8, fill: {{ color: C.white }}, line: {{ color: "D0E8F2", width: 1 }}, shadow: makeShadow() }});
          s.addShape(pres.shapes.RECTANGLE, {{ x: cx, y: 1.3, w: 0.08, h: 2.8, fill: {{ color: C.secondary }}, line: {{ color: C.secondary }} }});
          s.addText(String(b), {{ x: cx + 0.2, y: 1.5, w: cardW - 0.4, h: 2.4, fontSize: 14, fontFace: "Calibri", color: C.text_dark, align: "left", valign: "middle", wrap: true }});
        }});
    }}

    s.addShape(pres.shapes.RECTANGLE, {{ x: 0, y: 5.2, w: 10, h: 0.425, fill: {{ color: C.accent }}, line: {{ color: C.accent }} }});
    s.addText("Shiksha Sahayak", {{ x: 0.5, y: 5.2, w: 9, h: 0.425, fontSize: 9, color: C.white, align: "right", valign: "middle", margin: 0 }});

  }} else if (layout === "two_column") {{
    s.background = {{ color: C.white }};
    s.addShape(pres.shapes.RECTANGLE, {{ x: 0, y: 0, w: 10, h: 1.0, fill: {{ color: C.primary }}, line: {{ color: C.primary }} }});
    s.addText(safeTitle, {{ x: 0.5, y: 0, w: 9, h: 1.0, fontSize: 22, fontFace: "Cambria", bold: true, color: C.white, align: "left", valign: "middle", margin: 0 }});

    const half = Math.ceil(bullets.length / 2);
    const leftBullets = bullets.slice(0, half);
    const rightBullets = bullets.slice(half);

    s.addShape(pres.shapes.RECTANGLE, {{ x: 0.4, y: 1.2, w: 4.3, h: 4.0, fill: {{ color: C.light_bg }}, line: {{ color: "D0E8F2", width: 1 }}, shadow: makeShadow() }});
    if (leftBullets.length > 0) {{
        s.addText(leftBullets.map(b => ({{ text: String(b), options: {{ bullet: true, breakLine: true, paraSpaceAfter: 4 }} }})), {{ x: 0.6, y: 1.3, w: 4.0, h: 3.8, fontSize: 14, fontFace: "Calibri", color: C.text_dark, valign: "top" }});
    }}

    s.addShape(pres.shapes.RECTANGLE, {{ x: 5.2, y: 1.2, w: 4.3, h: 4.0, fill: {{ color: C.light_bg }}, line: {{ color: "D0E8F2", width: 1 }}, shadow: makeShadow() }});
    if (rightBullets.length > 0) {{
        s.addText(rightBullets.map(b => ({{ text: String(b), options: {{ bullet: true, breakLine: true, paraSpaceAfter: 4 }} }})), {{ x: 5.4, y: 1.3, w: 4.0, h: 3.8, fontSize: 14, fontFace: "Calibri", color: C.text_dark, valign: "top" }});
    }}

    s.addShape(pres.shapes.RECTANGLE, {{ x: 0, y: 5.2, w: 10, h: 0.425, fill: {{ color: C.accent }}, line: {{ color: C.accent }} }});
    s.addText("Shiksha Sahayak", {{ x: 0.5, y: 5.2, w: 9, h: 0.425, fontSize: 9, color: C.white, align: "right", valign: "middle", margin: 0 }});

  }} else {{
    s.background = {{ color: C.white }};
    s.addShape(pres.shapes.RECTANGLE, {{ x: 0, y: 0, w: 10, h: 1.0, fill: {{ color: C.primary }}, line: {{ color: C.primary }} }});
    s.addShape(pres.shapes.RECTANGLE, {{ x: 0, y: 0, w: 0.25, h: 1.0, fill: {{ color: C.secondary }}, line: {{ color: C.secondary }} }});
    s.addText(safeTitle, {{ x: 0.5, y: 0, w: 9, h: 1.0, fontSize: 22, fontFace: "Cambria", bold: true, color: C.white, align: "left", valign: "middle", margin: 0 }});
    s.addShape(pres.shapes.ROUNDED_RECTANGLE, {{ x: 9.0, y: 0.3, w: 0.5, h: 0.3, fill: {{ color: C.accent }}, line: {{ color: C.accent }}, rectRadius: 0.05 }});
    s.addText(String(idx + 1), {{ x: 9.0, y: 0.3, w: 0.5, h: 0.3, fontSize: 9, color: C.white, align: "center", valign: "middle", bold: true, margin: 0 }});
    s.addShape(pres.shapes.RECTANGLE, {{ x: 0.4, y: 1.15, w: 9.2, h: 4.05, fill: {{ color: C.light_bg }}, line: {{ color: "D0E8F2", width: 1 }}, shadow: makeShadow() }});

    if (bullets.length > 0) {{
        s.addText(bullets.map(b => ({{ text: String(b), options: {{ bullet: true, breakLine: true, paraSpaceAfter: 6 }} }})), {{ x: 0.65, y: 1.3, w: 8.7, h: 3.75, fontSize: 15, fontFace: "Calibri", color: C.text_dark, valign: "top" }});
    }}

    s.addShape(pres.shapes.RECTANGLE, {{ x: 0, y: 5.2, w: 10, h: 0.425, fill: {{ color: C.accent }}, line: {{ color: C.accent }} }});
    s.addText("Shiksha Sahayak", {{ x: 0.5, y: 5.2, w: 9, h: 0.425, fontSize: 9, color: C.white, align: "right", valign: "middle", margin: 0 }});
  }}
}});

pres.writeFile({{ fileName: "{output_path_js}" }})
  .then(() => console.log("PPT_DONE"))
  .catch(err => {{ console.error("PPT_ERROR:", err.message); process.exit(1); }});
"""

    with tempfile.NamedTemporaryFile(suffix=".js", dir=os.getcwd(), delete=False, mode="w", encoding="utf-8") as f:
        f.write(node_script)
        js_path = f.name

    try:
        result = subprocess.run(
            ["node", js_path],
            capture_output=True, text=True, timeout=60, encoding="utf-8"
        )
        if result.returncode != 0 or "PPT_ERROR" in result.stdout:
            raise RuntimeError(f"pptxgenjs failed: {result.stderr or result.stdout}")
        if "PPT_DONE" not in result.stdout:
            raise RuntimeError(f"pptxgenjs did not complete: {result.stdout}")
    finally:
        os.unlink(js_path)

def generate_ppt(file_storage, mode: str, slide_topics: list, slide_count: int, language: str = "English", grade: str = "Primary School") -> str:
    """
    Full pipeline: extract text (if provided) → Gemini generates slides → pptxgenjs builds PPTX.
    Returns the output file path.
    """
    text = ""
    if file_storage:
        text = extract_text_from_file(file_storage)
        if not text.strip():
            raise ValueError("No readable text found in the uploaded file.")

    # Pass the language and grade to Gemini
    slides = generate_slide_data_with_gemini(text, mode, slide_topics, slide_count, language, grade)

    output_dir = os.path.join(os.getcwd(), "generated_ppts")
    os.makedirs(output_dir, exist_ok=True)
    filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
    output_path = os.path.join(output_dir, filename)

    build_pptx_from_slides(slides, output_path)

    return output_path